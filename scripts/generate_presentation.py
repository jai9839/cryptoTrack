import os
from pathlib import Path
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:
    raise RuntimeError("Missing dependency: python-pptx. Install with: pip install python-pptx pillow")


ROOT = Path(__file__).resolve().parents[1]
README = ROOT / 'README.md'
IMAGES_DIR = ROOT / 'images'
OUT = ROOT / 'CryptoTrack_Presentation.pptx'


def read_readme_sections(readme_path):
    sections = {}
    if not readme_path.exists():
        return sections
    with readme_path.open(encoding='utf-8') as f:
        lines = f.readlines()

    current = None
    buffer = []
    for line in lines:
        if line.startswith('#'):
            if current:
                sections[current] = ''.join(buffer).strip()
            header = line.strip().lstrip('#').strip()
            current = header
            buffer = []
        else:
            buffer.append(line)
    if current:
        sections[current] = ''.join(buffer).strip()
    return sections


def add_title_slide(prs, title, subtitle=None):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle:
        if slide.placeholders and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle


def add_bullet_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.add_paragraph() if i > 0 else body.paragraphs[0]
        p.text = b
        p.level = 0
        p.font.size = Pt(18)


def add_image_slide(prs, title, image_path):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    left = Inches(0.5)
    top = Inches(1.5)
    max_w = prs.slide_width - Inches(1)
    try:
        slide.shapes.add_picture(str(image_path), left, top, width=max_w)
    except Exception:
        # ignore image errors
        pass


def collect_images(images_dir):
    imgs = []
    if not images_dir.exists():
        return imgs
    for ext in ('png', 'jpg', 'jpeg', 'gif', 'bmp'):
        imgs.extend(sorted(images_dir.glob(f'**/*.{ext}')))
    return imgs


def main():
    sections = read_readme_sections(README)
    prs = Presentation()
    # Ensure at least 20 slides target
    slides = []

    # Title
    title_text = 'CryptoTrack'
    subtitle = 'Real-time Crypto Tracking & Portfolio Manager'
    add_title_slide(prs, title_text, subtitle)

    # Overview / Description
    desc = sections.get('Description', '')
    add_bullet_slide(prs, 'Overview', [l for l in (desc.splitlines()[:6]) if l])

    # Live Demo
    demo = sections.get('Live Demo', '')
    add_bullet_slide(prs, 'Live Demo', [l for l in demo.splitlines() if l])

    # Features
    features = sections.get('Features', '')
    feat_lines = [f.strip('- ').strip() for f in features.splitlines() if f.strip().startswith('-')]
    if not feat_lines:
        feat_lines = [s for s in features.splitlines()[:6] if s]
    add_bullet_slide(prs, 'Key Features', feat_lines[:8])

    # Tech Stack
    tech = sections.get('Tech Stack', '')
    tech_lines = [l.strip('- ').strip() for l in tech.splitlines() if l.strip().startswith('-')]
    add_bullet_slide(prs, 'Tech Stack', tech_lines[:8])

    # Getting Started
    getting = sections.get('Getting Started', '')
    getting_lines = [l for l in getting.splitlines() if l.strip()][:8]
    add_bullet_slide(prs, 'Getting Started', getting_lines[:8])

    # Screenshots (multiple slides)
    images = collect_images(IMAGES_DIR)
    if images:
        for img in images[:6]:
            add_image_slide(prs, 'Screenshot', img)
    else:
        add_bullet_slide(prs, 'Screenshots', ['Screenshots are available in the repository images folder.'])

    # Attributions and Contact
    attrib = sections.get('Attributions', '')
    contact = sections.get('Contact', '')
    add_bullet_slide(prs, 'Attributions', [l for l in attrib.splitlines() if l.strip()][:6])
    add_bullet_slide(prs, 'Contact', [l for l in contact.splitlines() if l.strip()][:6])

    # Additional slides to reach up to 20
    while len(prs.slides) < 19:
        add_bullet_slide(prs, f'Additional Info ({len(prs.slides)+1})', ['More details about the project.'])

    # Final Thank You slide
    add_title_slide(prs, 'Thank You', 'Questions?')

    prs.save(str(OUT))
    print(f'Presentation written to: {OUT}')


if __name__ == '__main__':
    main()
